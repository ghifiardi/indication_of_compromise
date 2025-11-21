#!/usr/bin/env python3
"""
GATRA Executive Briefing Deck Generator
Creates a comprehensive PowerPoint presentation based on the white paper and demo
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_gatra_executive_deck():
    """Create comprehensive GATRA executive briefing deck"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    # Define colors
    dark_blue = RGBColor(0, 51, 102)
    gatra_green = RGBColor(0, 255, 65)  # Matrix green from demo
    red_alert = RGBColor(255, 68, 68)
    orange_warning = RGBColor(255, 165, 0)
    dark_bg = RGBColor(10, 10, 10)
    white = RGBColor(255, 255, 255)
    gray = RGBColor(128, 128, 128)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = dark_bg
    background.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "GATRA"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = white
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-DRIVEN SECURITY PLATFORM"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = gatra_green
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Revolutionizing SOC Operations with Autonomous Threat Detection"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(18)
    tagline_para.font.color.rgb = white
    tagline_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "The Challenge"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = red_alert
    
    p = tf.add_paragraph()
    p.text = "• Traditional SOCs take 200+ days to detect breaches"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• AI-driven attacks operate at machine speed (100+ requests/minute)"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• GTG-1002 campaign: 100M+ PII records targeted in <48 hours"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "The GATRA Solution"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = gatra_green
    
    p = tf.add_paragraph()
    p.text = "• Detection in 15-30 minutes (1,333x faster)"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 95% detection accuracy vs. <10% traditional"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Prevents $26.4M+ in breach costs per incident"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    # Slide 3: The GTG-1002 Threat
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The GTG-1002 Campaign: AI-Orchestrated Attack"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Six-Phase Attack Lifecycle"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    phases = [
        ("Phase 1: Initialization (45 seconds)", "Social engineering establishes access"),
        ("Phase 2: Reconnaissance (2-8 hours)", "150+ API requests/minute - CRITICAL DETECTION WINDOW"),
        ("Phase 3: Exploitation (1-4 hours)", "Legitimate tools exploited for persistence"),
        ("Phase 4: Lateral Movement (3-7 hours)", "Credential harvesting across 10+ systems"),
        ("Phase 5: Data Extraction (2-6 hours)", "100MB+ data exfiltration"),
        ("Phase 6: Handoff (Passive)", "Control passed to human operators")
    ]
    
    for i, (phase, description) in enumerate(phases):
        p = tf.add_paragraph()
        p.text = phase
        p.font.size = Pt(14)
        p.font.bold = True
        if i == 1:  # Phase 2 - critical
            p.font.color.rgb = red_alert
        
        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(12)
        p.level = 1
        if i == 1:
            p.font.color.rgb = red_alert
    
    # Slide 4: GATRA Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "GATRA: Five-Agent Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    agents = [
        ("Agent 1: Tempo Anomaly Detector", "Monitors API velocity (>150 req/min)", "88% confidence"),
        ("Agent 2: Phase Progression Tracker", "Recognizes kill-chain transitions", "88% confidence"),
        ("Agent 3: Data Exfiltration Analyzer", "Detects >100MB transfers", "98% confidence"),
        ("Agent 4: Lateral Movement Tracker", "Monitors credential pivoting", "87% confidence"),
        ("Agent 5: MCP Infrastructure Detector", "Identifies C2 patterns", "95% confidence")
    ]
    
    tf.text = "Multi-Agent Defense System"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = gatra_green
    
    for agent, function, performance in agents:
        p = tf.add_paragraph()
        p.text = agent
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = f"Function: {function}"
        p.font.size = Pt(12)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"Performance: {performance}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = gatra_green
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = ""
    
    # Slide 5: Live Demo Screenshot
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "GATRA Live Demo: Real-Time Threat Detection"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Interactive Simulation Features"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    demo_features = [
        "Real-time telemetry feed with live attack progression",
        "Visual agent status monitoring (4 agents active)",
        "Attack progression bar showing threat escalation",
        "Automated response demonstration",
        "15-second complete attack neutralization",
        "Zero human intervention required"
    ]
    
    for feature in demo_features:
        p = tf.add_paragraph()
        p.text = f"• {feature}"
        p.font.size = Pt(14)
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "Demo Results: GTG-1002 Attack Blocked in 15 Seconds"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = gatra_green
    
    # Slide 6: Comparative Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "GATRA vs. Traditional SOC Performance"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    # Create table manually using text boxes
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Performance Comparison"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    comparisons = [
        ("Detection Time", "15-30 minutes", "200+ days", "1,333x Faster"),
        ("Success Rate", "95%", "<10%", "+85% Efficiency"),
        ("False Positives", "<1%", "30%+", "95% Noise Reduction"),
        ("Response Type", "Autonomous", "Manual/Analyst", "Zero Human Delay"),
        ("Phase 2 Coverage", "88%", "5%", "Prevention Capability")
    ]
    
    for capability, gatra_val, traditional_val, advantage in comparisons:
        p = tf.add_paragraph()
        p.text = f"{capability}:"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = f"GATRA: {gatra_val} | Traditional: {traditional_val}"
        p.font.size = Pt(12)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"Advantage: {advantage}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = gatra_green
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = ""
    
    # Slide 7: ROI Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Return on Investment (ROI)"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Cost Avoidance Model"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = gatra_green
    
    p = tf.add_paragraph()
    p.text = "Single Prevented Breach Value:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    roi_items = [
        "Regulatory fines avoided: $15M+",
        "Remediation costs saved: $5M+",
        "Legal and notification costs: $3M+",
        "Brand damage prevention: $2M+",
        "Lost business prevention: $1.4M+",
        "Total breach cost avoided: $26.4M+"
    ]
    
    for item in roi_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "First Month ROI: 5,280%"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = gatra_green
    
    p = tf.add_paragraph()
    p.text = "Operational Efficiency: 80% SOC workload reduction"
    p.font.size = Pt(16)
    p.font.bold = True
    
    # Slide 8: Strategic Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Strategic Roadmap"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    roadmap_phases = [
        ("Current (2024-Q1 2025)", "Proven 5-Agent Solution", "• Operational and validated\n• 15-30 minute detection\n• GTG-1002 case study proven"),
        ("Near-Term (Q2-Q4 2025)", "Enhanced Platform", "• 10+ specialized agents\n• Industry-specific models\n• Compliance automation (GDPR, HIPAA)"),
        ("Vision (2026+)", "Global Intelligence Platform", "• Zero-day vulnerability detection\n• AI-vs-AI adversarial testing\n• Global threat intelligence network")
    ]
    
    for timeframe, phase_name, details in roadmap_phases:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        if tf.text:
            p.text = ""
        else:
            tf.text = ""
        
        p = tf.add_paragraph()
        p.text = f"{timeframe}: {phase_name}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = gatra_green
        
        p = tf.add_paragraph()
        p.text = details
        p.font.size = Pt(12)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = ""
    
    # Slide 9: Demo Call-to-Action
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Experience GATRA Live"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Interactive Demo Available"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = gatra_green
    
    p = tf.add_paragraph()
    p.text = "See GATRA detect and neutralize a GTG-1002 style attack in real-time:"
    p.font.size = Pt(16)
    
    demo_points = [
        "Watch all 5 agents coordinate in real-time",
        "See attack progression from breach to neutralization",
        "Experience 15-second threat resolution",
        "Understand the Phase 2 detection advantage",
        "Witness zero human intervention response"
    ]
    
    for point in demo_points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(14)
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "Demo File: demo.html (included in repository)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = gatra_green
    
    # Slide 10: Next Steps & Contact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Next Steps"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Recommended Actions"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    
    next_steps = [
        "1. Schedule live GTG-1002 simulation demo",
        "2. Review technical architecture documentation",
        "3. Discuss deployment timeline and pricing",
        "4. Evaluate integration with existing SOC infrastructure",
        "5. Plan pilot deployment strategy"
    ]
    
    for step in next_steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(16)
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "Contact Information"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = gatra_green
    
    p = tf.add_paragraph()
    p.text = "Repository: https://github.com/ghifiardi/indication_of_compromise"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Demo: Open demo.html in any web browser"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Documentation: Complete deployment guides included"
    p.font.size = Pt(14)
    p.level = 1
    
    # Save presentation
    output_file = "GATRA_Executive_Briefing_Deck.pptx"
    prs.save(output_file)
    print(f"✓ Executive briefing deck created: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    try:
        create_gatra_executive_deck()
    except ImportError:
        print("Error: python-pptx library not found.")
        print("Please install it using: pip install python-pptx")
    except Exception as e:
        print(f"Error creating presentation: {e}")
