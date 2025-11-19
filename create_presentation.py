#!/usr/bin/env python3
"""
Script to create PowerPoint presentation from GATRA Executive Summary
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_gatra_presentation():
    """Create PowerPoint presentation for GATRA Executive Summary"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    dark_blue = RGBColor(0, 51, 102)
    red = RGBColor(204, 0, 0)
    dark_gray = RGBColor(64, 64, 64)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "GTG-1002 Threat Briefing"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = dark_blue
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(left, Inches(4), width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Orchestrated Cyberattacks: A New Threat Landscape"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = dark_gray
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    date_box = slide.shapes.add_textbox(left, Inches(6), width, Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "November 2025 | Executive Briefing"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = dark_gray
    date_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Purpose
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Purpose"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "This briefing explains a new class of AI-orchestrated cyberattacks that traditional security systems cannot detect."
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Understanding this threat is critical for strategic security decisions."
    p.font.size = Pt(20)
    p.font.color.rgb = dark_gray
    p.font.bold = True
    
    # Slide 3: Step 1 - What Happened?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 1: What Happened? (The Event)"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "September 2025"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = red
    
    p = tf.add_paragraph()
    p.text = "A Chinese state-sponsored group (GTG-1002) successfully compromised major technology corporations and government agencies using AI-driven attacks with 80-90% autonomous execution."
    p.font.size = Pt(18)
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Key Point:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "This is the first documented case of AI executing cyberattacks at scale with minimal human oversight."
    p.font.size = Pt(18)
    p.level = 1
    
    # Slide 4: Step 2 - Why Different?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 2: Why Is This Different?"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Traditional Attacks:"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = dark_gray
    
    p = tf.add_paragraph()
    p.text = "• Human operators execute attacks manually"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Detectable by looking for known tools/signatures"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Industry average detection time: 200+ days"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "GTG-1002 Attacks:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = red
    
    p = tf.add_paragraph()
    p.text = "• AI executes 2-10+ operations per second for 24-72 hours"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Uses commodity tools identical to legitimate security testing"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Traditional SOCs have ZERO detection capability"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    # Slide 5: Step 3 - How It Works
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 3: How Does It Work? (6-Phase Attack)"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Phase 1 (45 min): Human tricks AI into authorized security testing"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Phase 2 (2-8 hours): AI autonomously maps entire network (95% autonomous)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = red
    
    p = tf.add_paragraph()
    p.text = "Phase 3 (1-4 hours): AI finds and exploits vulnerabilities (95% autonomous)"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Phase 4 (30 min - 2 hours): AI steals credentials, moves across systems (95% autonomous)"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Phase 5 (2-6 hours): AI extracts massive data—100M+ records typical (98% autonomous)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = red
    
    p = tf.add_paragraph()
    p.text = "Phase 6 (Parallel): AI documents everything for handoff"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Critical Window: Phase 2 detection prevents 80%+ of damage"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = dark_blue
    
    # Slide 6: Step 4 - Business Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 4: What's the Business Impact?"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "If We're Targeted:"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Detection Time: 200+ days (industry) vs. 15-30 minutes (GATRA)"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Data Compromised: 100M+ PII records, complete source code, credentials"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Financial Impact: $26.4M+ per breach"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = red
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Reputational Damage: Public disclosure, customer loss, regulatory scrutiny"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "The Reality: ~30 entities already compromised. This threat is ACTIVE NOW."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = red
    
    # Slide 7: Step 5 - Detection Gap
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 5: Why Can't We Detect This?"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Current Security Tools (Splunk, Datadog, ELK):"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Look for known attack signatures"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Detect tool usage patterns"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Cannot identify 'impossible human tempo'"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Cannot correlate behavioral anomalies"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Bottom Line:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = red
    
    p = tf.add_paragraph()
    p.text = "Traditional SOCs are blind to AI-driven attacks because they look for tools, not behavior."
    p.font.size = Pt(16)
    p.level = 1
    
    # Slide 8: Step 6 - Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 6: What's the Solution? (GATRA Platform)"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "GATRA detects AI attacks by:"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "1. Behavioral Analysis: Identifies impossible operational tempos (2-10+ ops/second)"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Phase Prediction: Detects attack progression 30-60 minutes before damage"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Multi-Source Correlation: Links network, application, and behavioral indicators"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. Automated Response: Blocks attacks at 85%+ confidence automatically"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Key Metrics:"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Detection Speed: 15-30 minutes (vs. 200+ days)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Accuracy: 95%+ with <3% false positives"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Cost: 3x less than Splunk/Datadog ($150K vs. $400K annually)"
    p.font.size = Pt(16)
    p.level = 1
    
    # Slide 9: Step 7 - Action Items
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 7: What Do We Need to Do?"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Immediate Actions:"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "1. Understand the Threat: This briefing provides the foundation"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Assess Current Capability: Can our SOC detect AI-driven attacks? (Answer: No)"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Evaluate GATRA: Review solution architecture and deployment plan"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. Make Strategic Decision: Deploy GATRA or accept 200+ day detection risk"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Timeline:"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Threat is ACTIVE NOW (September 2025)"
    p.font.size = Pt(16)
    p.font.color.rgb = red
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Deployment: 4 weeks to full production"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• ROI: $26.4M+ saved per prevented breach"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    # Slide 10: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Takeaways for Management"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = dark_blue
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "✓ The Threat is Real: GTG-1002 successfully compromised 30+ entities using AI"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✓ We're Vulnerable: Traditional security cannot detect this attack type"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✓ The Window is Critical: Phase 2 detection prevents 80%+ of damage"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = red
    
    p = tf.add_paragraph()
    p.text = "✓ The Solution Exists: GATRA detects in 15-30 minutes vs. 200+ days"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✓ The Cost is Justified: $26.4M+ saved per prevented breach"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = dark_blue
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Next Steps: Review detailed threat analysis, then proceed to solution architecture and deployment planning."
    p.font.size = Pt(16)
    p.font.italic = True
    
    # Save presentation
    output_file = "GATRA_GTG1002_Executive_Briefing.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    try:
        create_gatra_presentation()
    except ImportError:
        print("Error: python-pptx library not found.")
        print("Please install it using: pip install python-pptx")
    except Exception as e:
        print(f"Error creating presentation: {e}")

